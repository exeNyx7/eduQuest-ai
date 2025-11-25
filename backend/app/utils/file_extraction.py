"""
Utility functions for extracting text from various file formats.
"""
import base64
import re
from io import BytesIO
from typing import Optional

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def extract_text_from_content(content: str) -> str:
    """
    Extract text from content that might contain binary data markers.
    Supports: plain text, PDF (base64), PPTX (base64)
    Handles multiple files by extracting all binary markers.
    
    Args:
        content: Content string, potentially with [BINARY:filename:base64data] markers
        
    Returns:
        Extracted text content from all files combined
    """
    # Check if content contains binary data markers
    binary_pattern = r'\[BINARY:([^:]+):([^\]]+)\]'
    matches = re.findall(binary_pattern, content)
    
    if not matches:
        # Plain text content
        return content
    
    # Extract text from all binary files
    extracted_texts = []
    
    for filename, base64_data in matches:
        try:
            # Decode base64
            file_bytes = base64.b64decode(base64_data)
            
            # Determine file type and extract text
            if filename.lower().endswith('.pdf'):
                text = extract_text_from_pdf(file_bytes)
            elif filename.lower().endswith(('.pptx', '.ppt')):
                text = extract_text_from_pptx(file_bytes)
            else:
                text = f"[Unsupported file type: {filename}]"
            
            extracted_texts.append(f"=== Content from {filename} ===\n{text}")
            print(f"[FILE EXTRACTION] Extracted {len(text)} characters from {filename}")
            
        except Exception as e:
            print(f"[FILE EXTRACTION ERROR] Failed to extract from {filename}: {str(e)}")
            extracted_texts.append(f"[Error extracting text from {filename}: {str(e)}]")
    
    # Combine all extracted texts
    combined = "\n\n".join(extracted_texts)
    print(f"[FILE EXTRACTION] Total extracted content: {len(combined)} characters from {len(matches)} files")
    return combined


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """
    Extract text from PDF bytes.
    
    Args:
        file_bytes: PDF file content as bytes
        
    Returns:
        Extracted text
    """
    if PdfReader is None:
        raise ImportError("PyPDF2 is not installed. Install with: pip install PyPDF2")
    
    try:
        pdf_file = BytesIO(file_bytes)
        pdf_reader = PdfReader(pdf_file)
        
        text_parts = []
        for page in pdf_reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        
        full_text = "\n\n".join(text_parts)
        return full_text.strip() if full_text else "[No text found in PDF]"
        
    except Exception as e:
        print(f"[PDF EXTRACTION ERROR] {str(e)}")
        return f"[Error extracting PDF: {str(e)}]"


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """
    Extract text from PPTX bytes.
    
    Args:
        file_bytes: PPTX file content as bytes
        
    Returns:
        Extracted text
    """
    if Presentation is None:
        raise ImportError("python-pptx is not installed. Install with: pip install python-pptx")
    
    try:
        pptx_file = BytesIO(file_bytes)
        presentation = Presentation(pptx_file)
        
        text_parts = []
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_texts.append(shape.text)
            
            if slide_texts:
                text_parts.append(f"Slide {slide_num}:\n" + "\n".join(slide_texts))
        
        full_text = "\n\n".join(text_parts)
        return full_text.strip() if full_text else "[No text found in PPTX]"
        
    except Exception as e:
        print(f"[PPTX EXTRACTION ERROR] {str(e)}")
        return f"[Error extracting PPTX: {str(e)}]"
