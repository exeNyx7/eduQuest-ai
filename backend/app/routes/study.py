"""
Study Hall Routes - The Arcane Library
Manages user's uploaded study materials (Scrolls) and provides RAG-powered interactions
"""
from fastapi import APIRouter, HTTPException, UploadFile, File, Form
from pydantic import BaseModel
from typing import List, Optional, Dict
from datetime import datetime
import os
import io
import base64

from app.services.vector_store import store_content, retrieve_context
from app.services.ai_engine import AIEngine
from app.config.db import get_collection

# File parsing imports
try:
    from pypdf import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

router = APIRouter(tags=["study"])

# File parsing utilities
def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from PDF file"""
    if not PDF_AVAILABLE:
        raise HTTPException(status_code=400, detail="PDF parsing not available. Install pypdf.")
    
    try:
        pdf_file = io.BytesIO(file_bytes)
        reader = PdfReader(pdf_file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to parse PDF: {str(e)}")

def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text from Word document"""
    if not DOCX_AVAILABLE:
        raise HTTPException(status_code=400, detail="Word parsing not available. Install python-docx.")
    
    try:
        docx_file = io.BytesIO(file_bytes)
        doc = docx.Document(docx_file)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to parse Word document: {str(e)}")

def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from PowerPoint presentation"""
    if not PPTX_AVAILABLE:
        raise HTTPException(status_code=400, detail="PowerPoint parsing not available. Install python-pptx.")
    
    try:
        pptx_file = io.BytesIO(file_bytes)
        prs = Presentation(pptx_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to parse PowerPoint: {str(e)}")

def parse_file(file_bytes: bytes, filename: str) -> str:
    """Parse file based on extension"""
    ext = filename.lower().split('.')[-1]
    
    if ext == 'pdf':
        return extract_text_from_pdf(file_bytes)
    elif ext in ['docx', 'doc']:
        return extract_text_from_docx(file_bytes)
    elif ext in ['pptx', 'ppt']:
        return extract_text_from_pptx(file_bytes)
    elif ext in ['txt', 'md', 'text']:
        return file_bytes.decode('utf-8')
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {ext}")

# Pydantic Models
class UploadScrollRequest(BaseModel):
    """Request to upload a new Scroll (study material)"""
    user_id: str
    text: str
    filename: str
    file_type: str = "text"  # text, pdf, markdown, etc.
    topic: Optional[str] = None  # Optional topic/subject categorization


class ScrollMetadata(BaseModel):
    """Metadata for a stored Scroll"""
    scroll_id: str
    filename: str
    file_type: str
    topic: Optional[str]
    upload_date: str
    content_id: str
    chunks: int
    preview: str  # First 200 chars


class ChatWithScrollRequest(BaseModel):
    """Request to chat with a specific Scroll"""
    content_id: str
    user_query: str
    chat_history: Optional[List[Dict]] = []


class ChatResponse(BaseModel):
    """Response from chat"""
    response: str
    context_snippets: List[str]


class SummaryRequest(BaseModel):
    """Request to generate summary"""
    content_id: str
    max_length: str = "medium"  # short, medium, long


class SummaryResponse(BaseModel):
    """Summary response"""
    summary: str
    word_count: int


@router.post("/upload")
async def upload_scroll(req: UploadScrollRequest):
    """
    Upload a new Scroll to the Arcane Library.
    Stores vectors for RAG and metadata for library management.
    """
    try:
        print(f"[STUDY] Uploading scroll '{req.filename}' for user: {req.user_id}")
        
        # Store vectors (reuse existing logic)
        content_id = await store_content(req.user_id, req.text)
        print(f"[STUDY] Vectors stored with content_id: {content_id}")
        
        # Store metadata in study_materials collection
        materials_coll = get_collection("study_materials")
        
        # Calculate chunks (rough estimate based on text length)
        chunks = max(1, len(req.text.split()) // 100)
        
        # Create preview (first 200 chars)
        preview = req.text[:200].strip()
        if len(req.text) > 200:
            preview += "..."
        
        scroll_metadata = {
            "scroll_id": content_id,  # Use content_id as scroll_id for consistency
            "user_id": req.user_id,
            "filename": req.filename,
            "file_type": req.file_type,
            "topic": req.topic,
            "upload_date": datetime.utcnow().isoformat(),
            "content_id": content_id,
            "chunks": chunks,
            "preview": preview,
            "full_text": req.text  # Store full text for summaries
        }
        
        result = await materials_coll.insert_one(scroll_metadata)
        print(f"[STUDY] Metadata stored with _id: {result.inserted_id}")
        
        return {
            "success": True,
            "scroll_id": content_id,
            "content_id": content_id,
            "chunks": chunks,
            "message": f"📜 Scroll '{req.filename}' added to your Arcane Library!"
        }
        
    except Exception as e:
        print(f"[STUDY ERROR] Failed to upload scroll: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to upload scroll: {str(e)}")


@router.post("/upload-file")
async def upload_file_scroll(
    file: UploadFile = File(...),
    user_id: str = Form(...),
    topic: Optional[str] = Form(None)
):
    """
    Upload a file (PDF, Word, PowerPoint) to the Arcane Library.
    Parses the file and stores vectors for RAG.
    """
    try:
        print(f"[STUDY] Uploading file '{file.filename}' for user: {user_id}")
        
        # Read file bytes
        file_bytes = await file.read()
        
        # Parse file to extract text
        text = parse_file(file_bytes, file.filename)
        
        if not text or len(text.strip()) < 10:
            raise HTTPException(status_code=400, detail="File appears to be empty or text extraction failed")
        
        print(f"[STUDY] Extracted {len(text)} characters from {file.filename}")
        
        # Store vectors
        content_id = await store_content(user_id, text)
        print(f"[STUDY] Vectors stored with content_id: {content_id}")
        
        # Determine file type
        file_type = file.filename.lower().split('.')[-1]
        
        # Store metadata
        materials_coll = get_collection("study_materials")
        chunks = max(1, len(text.split()) // 100)
        preview = text[:200].strip()
        if len(text) > 200:
            preview += "..."
        
        scroll_metadata = {
            "scroll_id": content_id,
            "user_id": user_id,
            "filename": file.filename,
            "file_type": file_type,
            "topic": topic,
            "upload_date": datetime.utcnow().isoformat(),
            "content_id": content_id,
            "chunks": chunks,
            "preview": preview,
            "full_text": text
        }
        
        result = await materials_coll.insert_one(scroll_metadata)
        print(f"[STUDY] File metadata stored with _id: {result.inserted_id}")
        
        return {
            "success": True,
            "scroll_id": content_id,
            "content_id": content_id,
            "chunks": chunks,
            "message": f"📜 '{file.filename}' added to your Arcane Library!",
            "extracted_chars": len(text)
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"[STUDY ERROR] Failed to upload file: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to upload file: {str(e)}")


@router.get("/debug/{user_id}")
async def debug_user_scrolls(user_id: str):
    """Debug endpoint to check what's in the database"""
    try:
        from pymongo import MongoClient
        # Use pymongo (sync) to bypass any motor issues
        client = MongoClient('mongodb://localhost:27017/')
        db = client['eduquest']
        coll = db['study_materials']
        
        docs = list(coll.find({"user_id": user_id}).limit(10))
        
        return {
            "user_id": user_id,
            "found_docs": len(docs),
            "docs": [{
                "filename": doc.get("filename"),
                "user_id": doc.get("user_id"),
                "scroll_id": doc.get("scroll_id")
            } for doc in docs]
        }
    except Exception as e:
        return {"error": str(e)}


@router.get("/files/{user_id}")
async def get_user_scrolls(user_id: str):
    """
    Retrieve all Scrolls (study materials) for a user.
    Returns list of metadata for the user's library.
    """
    try:
        print(f"\n{'='*60}", flush=True)
        print(f"[STUDY] ⚡ GET /files/{user_id} called", flush=True)
        print(f"[STUDY] User ID: '{user_id}'", flush=True)
        print(f"[STUDY] User ID type: {type(user_id)}", flush=True)
        print(f"[STUDY] User ID length: {len(user_id)}", flush=True)
        print(f"{'='*60}\n", flush=True)
        
        # Use pymongo directly to bypass Motor issues
        try:
            from pymongo import MongoClient
            sync_client = MongoClient('mongodb://localhost:27017/')
            sync_db = sync_client['eduquest']
            sync_coll = sync_db['study_materials']
            
            # Check total count first
            total_count = sync_coll.count_documents({})
            print(f"[STUDY] Total documents in collection: {total_count}", flush=True)
            
            # Try query
            query = {"user_id": user_id}
            print(f"[STUDY] Query: {query}", flush=True)
            docs = list(sync_coll.find(query).sort("upload_date", -1).limit(100))
            print(f"[STUDY] ✅ Sync query found {len(docs)} documents", flush=True)
            
            scrolls = []
            for doc in docs:
                print(f"[STUDY] Processing doc: {doc.get('filename')}", flush=True)
                scrolls.append({
                    "scroll_id": str(doc["scroll_id"]),
                    "filename": doc["filename"],
                    "file_type": doc["file_type"],
                    "topic": doc.get("topic"),
                    "upload_date": doc["upload_date"],
                    "content_id": str(doc["content_id"]),
                    "chunks": doc["chunks"],
                    "preview": doc["preview"]
                })
            
            print(f"[STUDY] Returning {len(scrolls)} scrolls", flush=True)
            return {
                "scrolls": scrolls,
                "total": len(scrolls)
            }
            
        except Exception as sync_error:
            print(f"[STUDY] Sync fallback failed: {sync_error}", flush=True)
            # Return empty if both fail
            return {
                "scrolls": [],
                "total": 0
            }
        
    except Exception as e:
        print(f"[STUDY ERROR] Failed to fetch scrolls: {str(e)}", flush=True)
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Failed to fetch scrolls: {str(e)}")


@router.get("/content/{content_id}")
async def get_scroll_content(content_id: str, user_id: str):
    """
    Get full text content of a scroll by content_id.
    Used for quiz and flashcard generation.
    """
    try:
        print(f"[STUDY] Fetching content for content_id: {content_id}, user: {user_id}")
        
        from pymongo import MongoClient
        sync_client = MongoClient('mongodb://localhost:27017/')
        sync_db = sync_client['eduquest']
        sync_coll = sync_db['study_materials']
        
        doc = sync_coll.find_one({
            "content_id": content_id,
            "user_id": user_id
        })
        
        if not doc:
            print(f"[STUDY] ❌ Scroll not found for content_id: {content_id}")
            raise HTTPException(status_code=404, detail="Scroll not found")
        
        full_text = doc.get("full_text", "")
        print(f"[STUDY] ✅ Retrieved {len(full_text)} characters")
        
        return {
            "content_id": content_id,
            "filename": doc.get("filename"),
            "topic": doc.get("topic"),
            "text": full_text,
            "length": len(full_text)
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"[STUDY ERROR] Failed to fetch content: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch content: {str(e)}")


@router.post("/chat", response_model=ChatResponse)
async def chat_with_scroll(req: ChatWithScrollRequest):
    """
    Consult the Oracle - Chat with a specific Scroll using RAG.
    Retrieves relevant context and generates AI response.
    """
    try:
        print(f"[STUDY] Chat request for content_id: {req.content_id}")
        print(f"[STUDY] User query: {req.user_query[:100]}...")
        
        # Retrieve relevant context from vector store
        context_docs = await retrieve_context(req.user_query, req.content_id, limit=5)
        
        if not context_docs:
            return ChatResponse(
                response="I couldn't find relevant information in this Scroll to answer your question. Could you try rephrasing or ask about a different topic?",
                context_snippets=[]
            )
        
        context_texts = [d["text"] for d in context_docs]
        joined_context = "\n\n".join(context_texts)
        
        print(f"[STUDY] Retrieved {len(context_docs)} relevant chunks")
        
        # Generate AI response using chat_with_document method
        ai_engine = AIEngine()
        response = ai_engine.chat_with_document(
            context=joined_context,
            user_message=req.user_query,
            chat_history=req.chat_history
        )
        
        print(f"[STUDY] Generated response length: {len(response)} chars")
        
        return ChatResponse(
            response=response,
            context_snippets=context_texts[:3]  # Return top 3 snippets for reference
        )
        
    except Exception as e:
        print(f"[STUDY ERROR] Chat failed: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Chat failed: {str(e)}")


@router.post("/summary", response_model=SummaryResponse)
async def generate_scroll_summary(req: SummaryRequest):
    """
    Generate a Quick Recap summary of a Scroll.
    Retrieves full text and generates AI summary.
    """
    try:
        print(f"[STUDY] Summary request for content_id: {req.content_id}, length: {req.max_length}")
        
        # Retrieve full text from study_materials collection
        materials_coll = get_collection("study_materials")
        material = await materials_coll.find_one({"content_id": req.content_id})
        
        if not material:
            raise HTTPException(status_code=404, detail="Scroll not found")
        
        full_text = material.get("full_text", "")
        
        if not full_text:
            raise HTTPException(status_code=400, detail="No content available for summary")
        
        print(f"[STUDY] Generating {req.max_length} summary for text length: {len(full_text)} chars")
        
        # Generate summary using AI
        ai_engine = AIEngine()
        summary = ai_engine.generate_summary(full_text, max_length=req.max_length)
        
        word_count = len(summary.split())
        print(f"[STUDY] Generated summary with {word_count} words")
        
        return SummaryResponse(
            summary=summary,
            word_count=word_count
        )
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"[STUDY ERROR] Summary generation failed: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Summary generation failed: {str(e)}")


@router.delete("/files/{content_id}")
async def delete_scroll(content_id: str, user_id: str):
    """
    Delete a Scroll from the Arcane Library.
    Removes both metadata and vector data.
    """
    try:
        print(f"[STUDY] Deleting scroll {content_id} for user: {user_id}")
        
        # Delete metadata
        materials_coll = get_collection("study_materials")
        result = await materials_coll.delete_one({
            "content_id": content_id,
            "user_id": user_id
        })
        
        if result.deleted_count == 0:
            raise HTTPException(status_code=404, detail="Scroll not found or unauthorized")
        
        # Delete vector chunks
        docs_coll = get_collection("documents")
        vector_result = await docs_coll.delete_many({
            "content_id": content_id,
            "user_id": user_id
        })
        
        print(f"[STUDY] Deleted {result.deleted_count} metadata and {vector_result.deleted_count} vector chunks")
        
        return {
            "success": True,
            "message": "Scroll removed from your library",
            "deleted_chunks": vector_result.deleted_count
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"[STUDY ERROR] Failed to delete scroll: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to delete scroll: {str(e)}")
